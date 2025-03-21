from pptx import Presentation
from pptx.util import Inches
import os
import time
import warnings
warnings.filterwarnings("ignore", category=FutureWarning)

def create_report(image_path, result, filename):
    # Create a PowerPoint presentation object
    prs = Presentation()

    # Add title slide
    slide_layout = prs.slide_layouts[0]  # 0 is the title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Plastic Detection Report"
    subtitle.text = f"Analysis for {filename}"

    # Add slide for image
    slide_layout = prs.slide_layouts[5]  # 5 is a blank slide layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.add_picture(image_path, Inches(0.5), Inches(1), height=Inches(4.5))

    # Add slide for analysis result
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # 1 is the Title and Content layout
    title = slide.shapes.title
    title.text = "Analysis Result"

    content = slide.shapes.placeholders[1]
    content.text = result

    # Ensure the reports folder exists
    if not os.path.exists('reports'):
        os.makedirs('reports')

    # Create a unique filename for the report
    timestamp = int(time.time())
    report_filename = f"{timestamp}_{filename.split('.')[0]}_report.pptx"

    # Save the presentation to a file in the 'reports' folder
    report_path = os.path.join('reports', report_filename)
    prs.save(report_path)

    return report_path


def create_report_2(image_path, result, filename):
    # Create a PowerPoint presentation object
    prs = Presentation()

    # Add title slide
    slide_layout = prs.slide_layouts[0]  # 0 is the title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Plastic Detection Report"
    subtitle.text = f"Analysis for {filename}"

    # Add slide for image
    slide_layout = prs.slide_layouts[5]  # 5 is a blank slide layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.add_picture(image_path, Inches(0.5), Inches(1), height=Inches(4.5))

    # Add slide for analysis description
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Analysis Overview"
    content = slide.shapes.placeholders[1]
    content.text = result["description"]

    # Add slides for each object detected
    for obj in result['objects']:
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f"Object: {obj['type']}"

        content = slide.shapes.placeholders[1]
        content.text = (
            f"Shape (Aspect Ratio): {obj['shape']}\n"
            f"Color: {obj['color']}\n"
            f"Confidence: {obj['confidence']:.2f}"
        )

    # Ensure the reports folder exists
    if not os.path.exists('reports'):
        os.makedirs('reports')

    # Create a unique filename for the report
    timestamp = int(time.time())
    report_filename = f"{timestamp}_{filename.split('.')[0]}_report.pptx"

    # Save the presentation to a file in the 'reports' folder
    report_path = os.path.join('reports', report_filename)
    prs.save(report_path)

    return report_path

from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
import os
import time

def create_report_3(image_path, result, filename):
    # Ensure the reports folder exists
    if not os.path.exists('reports'):
        os.makedirs('reports')

    # Create a unique filename for the report
    timestamp = int(time.time())
    report_filename = f"{timestamp}_{filename.split('.')[0]}_report.pdf"
    report_path = os.path.join('reports', report_filename)

    # Create a canvas for the PDF
    c = canvas.Canvas(report_path, pagesize=letter)
    width, height = letter

    # Add title page
    c.setFont("Helvetica-Bold", 24)
    c.drawString(100, height - 100, "Plastic Detection Report")
    c.setFont("Helvetica", 12)
    c.drawString(100, height - 130, f"Analysis for {filename}")
    c.showPage()  # End the current page

    # Add image to PDF
    c.drawImage(image_path, 100, height - 450, width=400, height=300)
    c.showPage()  # End the current page

    # Add analysis description slide
    c.setFont("Helvetica-Bold", 18)
    c.drawString(100, height - 100, "Analysis Overview")
    c.setFont("Helvetica", 12)
    c.drawString(100, height - 130, result["description"])
    c.showPage()  # End the current page

    # Add slides for each object detected
    for obj in result['objects']:
        c.setFont("Helvetica-Bold", 18)
        c.drawString(100, height - 100, f"Object: {obj['type']}")
        c.setFont("Helvetica", 12)
        c.drawString(100, height - 130, f"Shape (Aspect Ratio): {obj['shape']}")
        c.drawString(100, height - 150, f"Color: {obj['color']}")
        c.drawString(100, height - 170, f"Confidence: {obj['confidence']:.2f}")
        c.showPage()  # End the current page

    # Save the PDF file
    c.save()

    return report_path

def check_for_pyrolysis(object_type):
    """
    Check if the detected object is useful for the pyrolysis process based on its type.
    
    Args:
    - object_type (str): The type of the detected object (e.g., "plastic", "paper", "metal").
    
    Returns:
    - bool: True if the object is useful for pyrolysis, False otherwise.
    """
    # List of object types useful for pyrolysis (plastic types commonly used for pyrolysis)
    pyrolysis_compatible_types = [
        "polyethylene", "polypropylene", "polystyrene", "polyvinyl chloride", "polyethylene terephthalate",
        "bottle", "cup", "cover", "wrapper", "bag", "can", "container", "cell phone", "remote", "toilet","chair"
    ]

    # If the object type matches one of the compatible types, return True (useful for pyrolysis)
    if object_type.lower() in pyrolysis_compatible_types:
        return True
    else:
        return False


from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
import os
import time
from reportlab.lib import colors

def create_report_5(image_path, result, filename):
    # Ensure the reports folder exists
    if not os.path.exists('reports'):
        os.makedirs('reports')

    # Create a unique filename for the report
    timestamp = int(time.time())
    report_filename = f"{timestamp}_{filename.split('.')[0]}_report.pdf"
    report_path = os.path.join('reports', report_filename)

    # Create a canvas for the PDF
    c = canvas.Canvas(report_path, pagesize=letter)
    width, height = letter

    # **First page**: Display logo (Favicon.png)
    c.drawImage(r"static/images/Favicon.png", 100, height - 150, width=200, height=150)
    c.setFont("Helvetica-Bold", 24)
    c.drawString(100, height - 250, "Plastic Detection Report")
    c.setFont("Helvetica", 12)
    c.drawString(100, height - 280, f"Analysis for {filename}")
    c.showPage()  # End first page

    # **Second page**: Display the uploaded image
    c.setFont("Helvetica-Bold", 18)
    c.drawString(100, height - 100, "This is the image we have analyzed:")
    c.drawImage(image_path, 100, height - 450, width=400, height=300)
    c.showPage()  # End second page

    # **Third page**: Number of objects and their names
    c.setFont("Helvetica-Bold", 18)
    c.drawString(100, height - 100, "Number of Objects Detected and Their Names:")
    y_position = height - 130  # Starting position for list
    for obj in result['objects']:
        c.setFont("Helvetica", 12)
        c.drawString(100, y_position, f"- {obj['type']}")
        y_position -= 20  # Adjust the y-position for each object
    c.showPage()  # End third page

    # **From the 4th page onwards**: Description of each object
    for obj in result['objects']:
        c.setFont("Helvetica-Bold", 18)
        c.drawString(100, height - 100, f"Object: {obj['type']}")

        # # Display the object image (small size, 1/5th of the page size)
        # small_image_path = obj['image']  # You can provide the path to a cropped or resized image
        # c.drawImage(small_image_path, width - 200, height - 150, width=100, height=100)

        # Display object attributes
        c.setFont("Helvetica", 12)
        c.drawString(100, height - 180, f"Shape: {obj['shape']}")
        c.drawString(100, height - 200, f"Color: {obj['color']}")
        c.drawString(100, height - 220, f"Confidence: {obj['confidence']:.2f}")
        #c.drawString(100, height - 240, f"Useful for Pyrolysis: {'Yes' if obj['useful_for_pyrolysis'] else 'No'}")
        useful_for_pyrolysis = check_for_pyrolysis(obj['type'])
        c.drawString(100, height - 240, f"Useful for Pyrolysis: {'Yes' if useful_for_pyrolysis else 'No'}")


        # Check if there are more attributes to display
        if 'additional_attributes' in obj:
            y_position = height - 260
            for attribute, value in obj['additional_attributes'].items():
                c.drawString(100, y_position, f"{attribute}: {value}")
                y_position -= 20

        c.showPage()  # End page for the object

    # Save the PDF file
    c.save()

    return report_path


from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
import os
import time

# Define a function to convert RGB to color names (a basic approach)
def rgb_to_color_name(rgb):
    # You can extend this function based on your requirements
    color_names = {
        (255, 0, 0): 'Red',
        (0, 255, 0): 'Green',
        (0, 0, 255): 'Blue',
        (255, 255, 0): 'Yellow',
        (0, 255, 255): 'Cyan',
        (255, 0, 255): 'Magenta',
        (0, 0, 0): 'Black',
        (255, 255, 255): 'White',
    }
    return color_names.get(tuple(rgb), 'Unknown Color')  # Default to 'Unknown Color'

def create_report_4(image_path, result, filename):
    # Ensure the reports folder exists
    if not os.path.exists('reports'):
        os.makedirs('reports')

    # Create a unique filename for the report
    timestamp = int(time.time())
    file_base_name = filename.split('.')[0]  # Get the name without extension
    report_filename = f"plastic_detection_report_{timestamp}.pdf"
    report_path = os.path.join('reports', report_filename)

    # Create a canvas for the PDF
    c = canvas.Canvas(report_path, pagesize=letter)
    width, height = letter

    def draw_header_footer(page_num):
        # Draw header
        c.setFont("Helvetica-Bold", 10)
        c.drawString(100, height - 20, "Plastic Detection Report - Image Analysis")
        
        # Draw footer
        c.setFont("Helvetica", 10)  # Increased font size for better readability
        contact_text = "Contact the owner"
        c.drawString(100, 40, contact_text)  # Adjusted y-position to prevent overlap
        c.linkURL("https://www.linkedin.com/in/shaik-sameer-hussain-b88323250", (100, 30, 400, 50))

    # First page: Display logo (Favicon.png)
    draw_header_footer(1)
    c.drawImage(r"static/images/Favicon.png", 100, height - 150, width=200, height=150)
    c.setFont("Helvetica-Bold", 24)
    c.drawString(100, height - 250, "Plastic Detection Report")
    c.setFont("Helvetica", 12)
    c.drawString(100, height - 280, f"Analysis for {filename}")
    c.showPage()  # End first page

    # Second page: Display the uploaded image
    draw_header_footer(2)
    c.setFont("Helvetica-Bold", 18)
    c.drawString(100, height - 100, "This is the image we have analyzed:")
    c.drawImage(image_path, 100, height - 450, width=400, height=300)
    c.showPage()  # End second page

    # Third page: Number of objects and their names
    draw_header_footer(3)
    c.setFont("Helvetica-Bold", 18)
    c.drawString(100, height - 100, "Number of Objects Detected and Their Names:")
    y_position = height - 130  # Starting position for list
    for obj in result['objects']:
        c.setFont("Helvetica", 12)
        c.drawString(100, y_position, f"- {obj['type']}")
        y_position -= 20  # Adjust the y-position for each object
    c.showPage()  # End third page

    # From the 4th page onwards: Description of each object
    for idx, obj in enumerate(result['objects']):
        draw_header_footer(4 + idx)
        c.setFont("Helvetica-Bold", 18)
        c.drawString(100, height - 100, f"Object: {obj['type']}")

        # Display object attributes
        c.setFont("Helvetica", 12)
        c.drawString(100, height - 180, f"Shape: {obj['shape']}")
        
        # Convert RGB to color name
        color_name = rgb_to_color_name(obj['color'])  # Assuming color is in RGB format
        c.drawString(100, height - 200, f"Color: {color_name}")
        
        c.drawString(100, height - 220, f"Confidence: {obj['confidence']:.2f}")

        useful_for_pyrolysis = check_for_pyrolysis(obj['type'])
        c.drawString(100, height - 240, f"Useful for Pyrolysis: {'Yes' if useful_for_pyrolysis else 'No'}")

        # Display additional attributes if available
        if 'additional_attributes' in obj:
            y_position = height - 260
            for attribute, value in obj['additional_attributes'].items():
                c.drawString(100, y_position, f"{attribute}: {value}")
                y_position -= 20

        # Add dynamic closing statement for the object
        c.setFont("Helvetica", 10)
        c.drawString(100, height - 280, f"Conclusion: The detected object is analyzed based on its physical and chemical properties, providing insights into its composition and utility.")

        c.showPage()  # End page for the object

    # Last page: Contact information
    draw_header_footer("Final")
    c.setFont("Helvetica-Bold", 14)
    c.drawString(100, height - 100, "Thank you for using our Image Analysis service!")
    c.setFont("Helvetica", 12)
    c.drawString(100, height - 130, "If you have any doubts or inquiries, feel free to contact us:")

    # Link to LinkedIn profile
    c.setFont("Helvetica-Bold", 12)
    contact_text = "Contact: Click here to visit LinkedIn"
    c.drawString(100, height - 160, contact_text)
    # Make the contact text clickable and link to LinkedIn profile
    c.linkURL("https://www.linkedin.com/in/shaik-sameer-hussain-b88323250", (100, height - 170, 400, height - 150))

    c.showPage()  # End last page

    # Save the PDF file
    c.save()

    return report_path
